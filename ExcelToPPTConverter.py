import pandas as pd
from pptx import Presentation
from pptx.util import Inches

# Read and store content 
# of an excel file  
read_file = pd.read_excel ("Test.xlsx") 
  
# read csv file and convert  
# into a dataframe object 
df = pd.DataFrame(pd.read_excel("Test.xlsx")) 
  
col = df.columns
  
rowCount=df.shape[0]

root = Presentation()

for x in range(rowCount):

    first_slide_layout = root.slide_layouts[6]
    slide = root.slides.add_slide(first_slide_layout)


    rows = len(col)+1
    cols = 2
    left = top = Inches(0.2)
    width = Inches(6.0)
    height = Inches(1.0)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    

    # set column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(4.0)

    
    for y in range(1,len(col)+1):
        table.cell(y, 0).text = df.columns.astype(str).values[y-1]
        table.cell(y, 1).text = df.loc[x].astype(str).values[y-1]

root.save("ConvertedPPTX.pptx")





